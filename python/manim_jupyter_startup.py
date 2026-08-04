"""Private kernel runtime for ``*.manim.ipynb`` notebooks.

The VS Code extension injects this module into its private Jupyter kernel. The
TypeScript controller decides whether a Cell is Python or Manim and wraps only
explicit Manim Cells before execution.
"""

import logging as _manim_jupyter_logging
import mimetypes as _manim_jupyter_mimetypes
from pathlib import Path as _ManimJupyterPath

from IPython.display import clear_output as _manim_jupyter_clear_output
from manim import *
import manim.utils.ipython_magic as _manim_jupyter_ipython_magic
from manim_slides import Slide as _ManimJupyterSlide
from manim_slides import ThreeDSlide as _ManimJupyterThreeDSlide
from manim_slides.config import BaseSlideConfig as _ManimJupyterBaseSlideConfig


_manim_jupyter_options = globals().get("_MANIM_JUPYTER_BOOTSTRAP", {})
_MANIM_JUPYTER_CELL_SETTINGS = _manim_jupyter_options.get("cellSettings", {})
_MANIM_JUPYTER_ACTIVE_CELL_SETTINGS = {}
_MANIM_JUPYTER_MAGIC_ARGS = _manim_jupyter_options.get(
    "magicArgs", "-ql -v WARNING --progress_bar display {scene}"
)
MANIM_THEME = _manim_jupyter_options.get("theme", "dark")
MANIM_FOREGROUND = _manim_jupyter_options.get("foregroundColor", "#F8FAFC")

config.media_width = _manim_jupyter_options.get("mediaWidth", "100%")
# The companion renderer reads the generated file in bounded chunks. Never
# turn a video into a base64 HTML string: that multiplies its memory use across
# the Python kernel, extension host, notebook model, and output webview.
config.media_embed = False
config.progress_bar = "display"
config.background_color = _manim_jupyter_options.get("backgroundColor", "#0E1117")
config.frame_width = config.frame_height * float(_manim_jupyter_options.get("aspect", 16 / 9))
_manim_jupyter_logging.getLogger("manim-slides").setLevel(_manim_jupyter_logging.WARNING)

_MANIM_JUPYTER_VIDEO_MIME = "application/vnd.manim.video+json"
_ManimJupyterOriginalVideo = globals().get(
    "_ManimJupyterOriginalVideo", _manim_jupyter_ipython_magic.Video
)


class _ManimJupyterBoundedVideo(_ManimJupyterOriginalVideo):
    """Expose a file descriptor instead of an in-memory base64 video."""

    def __init__(self, data=None, url=None, filename=None, embed=False, *args, **kwargs):
        super().__init__(
            data=data,
            url=url,
            filename=filename,
            embed=False,
            *args,
            **kwargs,
        )

    def _repr_mimebundle_(self, include=None, exclude=None):
        if self.filename:
            # Keep rendering progress visible while Manim is working, then
            # atomically replace it with the only persistent output: the video.
            _manim_jupyter_clear_output(wait=True)
            candidate = _ManimJupyterPath(self.filename).resolve()
            scene_names = config.get("scene_names") or []
            scene_name = scene_names[0] if scene_names else ""
            options = _MANIM_JUPYTER_ACTIVE_CELL_SETTINGS or _MANIM_JUPYTER_CELL_SETTINGS.get(scene_name, {})
            mime_type = self.mimetype or _manim_jupyter_mimetypes.guess_type(candidate)[0]
            return {
                _MANIM_JUPYTER_VIDEO_MIME: {
                    "path": str(candidate),
                    "mimeType": mime_type or "video/mp4",
                    "loop": bool(_manim_jupyter_options.get("videoLoop", False)),
                    "controls": bool(options.get("controls", True)),
                    "playbackRate": float(options.get("playbackRate", 1.0)),
                    "width": str(config.media_width or "100%"),
                },
                "text/plain": f"Manim video: {candidate.name}",
            }
        return {"text/html": super()._repr_html_()}


_manim_jupyter_ipython_magic.Video = _ManimJupyterBoundedVideo

_ManimJupyterManimScene = Scene
_ManimJupyterManimThreeDScene = ThreeDScene
Scene = _ManimJupyterSlide
ThreeDScene = _ManimJupyterThreeDSlide
Slide = _ManimJupyterSlide
ThreeDSlide = _ManimJupyterThreeDSlide


_manim_jupyter_ip = get_ipython()
_manim_jupyter_ip.register_magics(_manim_jupyter_ipython_magic.ManimMagic)


_MANIM_JUPYTER_TIMING_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"

_ManimJupyterTimingXml = """<p:timing xmlns:p="{ns}"><p:tnLst><p:par><p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot"><p:childTnLst><p:seq concurrent="1" nextAc="seek"><p:cTn id="2" dur="indefinite" nodeType="mainSeq"><p:childTnLst><p:par><p:cTn id="3" fill="hold"><p:stCondLst><p:cond delay="indefinite"/><p:cond evt="onBegin" delay="0"><p:tn val="2"/></p:cond></p:stCondLst><p:childTnLst><p:par><p:cTn id="4" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst><p:childTnLst><p:par><p:cTn id="5" presetID="1" presetClass="mediacall" presetSubtype="0" fill="hold" nodeType="afterEffect"><p:stCondLst><p:cond delay="0"/></p:stCondLst><p:childTnLst><p:cmd type="call" cmd="playFrom(0.0)"><p:cBhvr><p:cTn id="6" dur="2000" fill="hold"/><p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl></p:cBhvr></p:cmd></p:childTnLst></p:cTn></p:par></p:childTnLst></p:cTn></p:par></p:childTnLst></p:cTn></p:par></p:childTnLst></p:cTn><p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst><p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst></p:seq><p:video><p:cMediaNode vol="80000"><p:cTn id="7" fill="hold" display="0"><p:stCondLst><p:cond delay="indefinite"/></p:stCondLst></p:cTn><p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl></p:cMediaNode></p:video><p:seq concurrent="1" nextAc="seek"><p:cTn id="8" restart="whenNotActive" fill="hold" evtFilter="cancelBubble" nodeType="interactiveSeq"><p:stCondLst><p:cond evt="onClick" delay="0"><p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl></p:cond></p:stCondLst><p:endSync evt="end" delay="0"><p:rtn val="all"/></p:endSync><p:childTnLst><p:par><p:cTn id="9" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst><p:childTnLst><p:par><p:cTn id="10" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst><p:childTnLst><p:par><p:cTn id="11" presetID="2" presetClass="mediacall" presetSubtype="0" fill="hold" nodeType="clickEffect"><p:stCondLst><p:cond delay="0"/></p:stCondLst><p:childTnLst><p:cmd type="call" cmd="togglePause"><p:cBhvr><p:cTn id="12" dur="1" fill="hold"/><p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl></p:cBhvr></p:cmd></p:childTnLst></p:cTn></p:par></p:childTnLst></p:cTn></p:par></p:childTnLst></p:cTn></p:par></p:childTnLst></p:cTn><p:nextCondLst><p:cond evt="onClick" delay="0"><p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl></p:cond></p:nextCondLst></p:seq></p:childTnLst></p:cTn></p:par></p:tnLst></p:timing>"""


def _ManimJupyterAutoPlayMedia(media, loop=False):
    """Make a python-pptx movie start automatically when its slide is shown.

    PowerPoint does not honour python-pptx's bare timing tree (a lone
    ``<p:video><p:cMediaNode delay="0">`` never starts the movie).  It needs
    the full native structure: a ``mainSeq`` whose first child triggers
    ``playFrom(0.0)`` on ``onBegin``, plus the media node and a click handler.
    The whole ``<p:timing>`` element is therefore replaced with exactly the
    tree PowerPoint itself writes for a slide with an autoplaying video.
    """
    import lxml.etree as _manim_jupyter_etree

    media_id = _manim_jupyter_etree.ElementBase.xpath(
        media.element,
        ".//p:cNvPr",
        namespaces={"p": _MANIM_JUPYTER_TIMING_NS},
    )[0].attrib["id"]
    slide = media.element.getparent().getparent().getparent()
    timing = _manim_jupyter_etree.ElementBase.xpath(
        slide,
        "./p:cSld/../p:timing",
        namespaces={"p": _MANIM_JUPYTER_TIMING_NS},
    )
    if not timing:
        return
    timing = timing[0]
    parent = timing.getparent()
    if parent is None:
        return
    parent.remove(timing)

    replacement = _manim_jupyter_etree.fromstring(
        _ManimJupyterTimingXml.format(
            ns=_MANIM_JUPYTER_TIMING_NS,
            spid=media_id,
        )
    )
    if loop:
        for node in _manim_jupyter_etree.ElementBase.xpath(
            replacement,
            ".//p:cTn[@id='5']|.//p:cTn[@id='7']",
            namespaces={"p": _MANIM_JUPYTER_TIMING_NS},
        ):
            node.set("repeatCount", "indefinite")
    parent.append(replacement)


def _ManimJupyterBuildPptx(video_files, destination, loop=False):
    """Build one PowerPoint slide per rendered Manim animation.

    python-pptx adds the video with ``delay="indefinite"`` (starts on click).
    This rewrite removes the wait condition and marks the media node
    ``visible="true"`` so the video starts playing automatically when the
    slide is shown, in every PowerPoint and LibreOffice version.
    """
    import mimetypes as _manim_jupyter_mimetypes
    import tempfile as _manim_jupyter_tempfile

    from manim_slides.convert import (
        FrameIndex as _ManimJupyterFrameIndex,
        read_image_from_video_file as _ManimJupyterReadFrame,
    )
    from pptx import Presentation as _ManimJupyterPresentation

    if not video_files:
        raise ValueError("没有可导出的 Manim 动画；请至少使用一次 self.play(...) 或 self.wait(...)。")

    prs = _ManimJupyterPresentation()
    resolution = globals().get("_MANIM_JUPYTER_PPTX_RESOLUTION")
    if not (
        isinstance(resolution, (tuple, list))
        and len(resolution) == 2
        and all(isinstance(value, int) and value > 0 for value in resolution)
    ):
        resolution = (
            int(config.get("pixel_width") or 1280),
            int(config.get("pixel_height") or 720),
        )
    pixel_width, pixel_height = resolution
    prs.slide_width = pixel_width * 9525
    prs.slide_height = pixel_height * 9525
    layout = prs.slide_layouts[6]  # blank layout

    with _manim_jupyter_tempfile.TemporaryDirectory() as directory_name:
        directory = _ManimJupyterPath(directory_name)
        for index, video_file in enumerate(video_files):
            video_path = _ManimJupyterPath(video_file)
            if not video_path.is_file():
                raise ValueError(f"Manim 动画视频不存在：{video_path}")

            poster_path = directory / f"{index:05d}.png"
            _ManimJupyterReadFrame(video_path, _ManimJupyterFrameIndex.first).save(
                poster_path
            )

            slide = prs.slides.add_slide(layout)
            mime_type = (
                _manim_jupyter_mimetypes.guess_type(video_path)[0] or "video/mp4"
            )
            movie = slide.shapes.add_movie(
                str(video_path),
                0,
                0,
                prs.slide_width,
                prs.slide_height,
                poster_frame_image=str(poster_path),
                mime_type=mime_type,
            )
            _ManimJupyterAutoPlayMedia(movie, loop=loop)

    destination = _ManimJupyterPath(destination)
    destination.parent.mkdir(parents=True, exist_ok=True)
    prs.save(destination)


_MANIM_JUPYTER_READY = True
